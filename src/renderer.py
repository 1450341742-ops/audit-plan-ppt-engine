from __future__ import annotations
from copy import deepcopy
from pathlib import Path
from typing import Any
import re
import math

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

BASE_DIR = Path(__file__).resolve().parent.parent
DEFAULT_TEMPLATE_PATH = BASE_DIR / "assets" / "template.pptx"

SLIDE_COVER = 2
SLIDE_THANKS = 3
SLIDE_TOC = 4
SLIDE_PART1 = 5
SLIDE_OVERVIEW = 6
SLIDE_SCOPE = 7
SLIDE_COMMON = 8
SLIDE_INDIVIDUAL = 9
SLIDE_PART2 = 10
SLIDE_COUNTS = 11
SLIDE_ISSUE = 12

WHITE = RGBColor(255, 255, 255)
YELLOW = RGBColor(255, 242, 0)
BLACK = RGBColor(0, 0, 0)
BLUE = RGBColor(0, 112, 192)
DARK_BLUE = RGBColor(31, 78, 121)
LIGHT_BLUE = RGBColor(221, 235, 247)
LIGHT_YELLOW = RGBColor(255, 242, 204)
GRAY = RGBColor(89, 89, 89)

LEFT_CATS = ["国家药物临床试验政策法规的遵循", "伦理委员会审核要求的遵循", "知情同意书（ICF）的签署和记录", "原始文件的建立、内容和记录", "门诊/住院HIS、LIS、PACS等系统数据溯源", "方案依从性", "药物疗效/研究评价指标的评估"]
RIGHT_CATS = ["安全性信息评估，记录与报告", "CRF填写（时效性、一致性、溯源性、完整性）", "试验用药品管理", "生物样本管理", "临床研究必须文件", "申办方/CRO职责", "其他"]
PLACEHOLDER_TEXTS = ["单击此处编辑标题", "单击此处编辑副标题", "Click to edit Master title style", "Click to edit Master subtitle style"]

RISK_CATEGORY_WEIGHT = {
    "知情同意书（ICF）的签署和记录": 45,
    "安全性信息评估，记录与报告": 45,
    "方案依从性": 42,
    "原始文件的建立、内容和记录": 38,
    "CRF填写（时效性、一致性、溯源性、完整性）": 36,
    "试验用药品管理": 36,
    "生物样本管理": 32,
    "伦理委员会审核要求的遵循": 32,
    "门诊/住院HIS、LIS、PACS等系统数据溯源": 30,
    "药物疗效/研究评价指标的评估": 30,
    "临床研究必须文件": 24,
    "国家药物临床试验政策法规的遵循": 22,
    "申办方/CRO职责": 20,
    "其他": 10,
}

RISK_KEYWORDS = {
    "SAE": 28, "SUSAR": 28, "死亡": 28, "住院": 22, "严重不良": 24, "安全性": 18,
    "知情同意": 26, "签署": 16, "受试者权益": 22, "伦理": 18,
    "入排": 24, "入组": 16, "排除标准": 24, "方案偏离": 22, "违背方案": 22,
    "主要终点": 20, "疗效评估": 18, "RECIST": 18,
    "原始记录": 20, "源文件": 20, "溯源": 18, "不一致": 18, "缺失": 16,
    "EDC": 18, "CRF": 18, "迟录": 14, "漏录": 16,
    "试验用药": 18, "药品": 16, "超温": 18, "发放": 14, "回收": 14,
    "样本": 16, "中心实验室": 14, "PACS": 12, "LIS": 12, "HIS": 12,
}


def _clean(v: Any) -> str:
    s = str(v or "").replace("\r", "\n")
    s = re.sub(r"[ \t\xa0]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def _shorten(text: str, limit: int = 90) -> str:
    text = _clean(text).replace("\n", "；")
    if len(text) <= limit:
        return text or "—"
    return text[:limit].rstrip("；，,。 ") + "…"


def _delete_slide(prs: Presentation, index: int = 0) -> None:
    r_id = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[index]


def _remove_shape_xml(shape) -> None:
    try:
        shape.element.getparent().remove(shape.element)
    except Exception:
        pass


def _text(shape) -> str:
    if getattr(shape, "has_text_frame", False):
        return _clean(shape.text)
    return ""


def _is_placeholder_text(text: str) -> bool:
    return any(p in _clean(text) for p in PLACEHOLDER_TEXTS)


def _copy_rels(src_part, dest_part) -> dict[str, str]:
    rel_map = {}
    for r_id, rel in src_part.rels.items():
        if "slideLayout" in rel.reltype:
            continue
        try:
            if rel.is_external:
                new_rid = dest_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            else:
                new_rid = dest_part.rels.get_or_add(rel.reltype, rel._target)
            rel_map[r_id] = new_rid
        except Exception:
            pass
    return rel_map


def _remap_relationship_ids(xml_el, rel_map: dict[str, str]) -> None:
    for el in xml_el.iter():
        for attr, val in list(el.attrib.items()):
            if val in rel_map:
                el.attrib[attr] = rel_map[val]


def _copy_slide(prs: Presentation, src_no: int):
    source = prs.slides[src_no - 1]
    dest = prs.slides.add_slide(source.slide_layout)
    for shp in list(dest.shapes):
        _remove_shape_xml(shp)
    rel_map = _copy_rels(source.part, dest.part)
    for shape in source.shapes:
        try:
            if _is_placeholder_text(_text(shape)):
                continue
            new_el = deepcopy(shape.element)
            _remap_relationship_ids(new_el, rel_map)
            dest.shapes._spTree.insert_element_before(new_el, "p:extLst")
        except Exception:
            pass
    return dest


def _remove_original_template_slides(prs: Presentation, original_count: int) -> None:
    for _ in range(original_count):
        _delete_slide(prs, 0)


def _tables(slide):
    return [shp.table for shp in slide.shapes if getattr(shp, "has_table", False)]


def _slide_text(slide) -> str:
    return "\n".join(_text(shp) for shp in slide.shapes if getattr(shp, "has_text_frame", False))


def _has_meaningful_non_text_content(slide) -> bool:
    for shp in slide.shapes:
        try:
            if getattr(shp, "has_table", False):
                return True
            if getattr(shp, "has_text_frame", False):
                if not _is_placeholder_text(_text(shp)) and _clean(_text(shp)):
                    return True
                continue
            return True
        except Exception:
            continue
    return False


def _remove_template_empty_slides(prs: Presentation) -> None:
    for idx in range(len(prs.slides) - 1, -1, -1):
        slide = prs.slides[idx]
        text = _slide_text(slide)
        if _is_placeholder_text(text) and not _has_meaningful_non_text_content(slide):
            _delete_slide(prs, idx)
            continue
        meaningful = re.sub(r"[\s\-—_：:|]+", "", text)
        if not meaningful and not _has_meaningful_non_text_content(slide):
            _delete_slide(prs, idx)


def _remove_text_shapes(slide, keep_keywords: tuple[str, ...] = ()) -> None:
    for shp in list(slide.shapes):
        try:
            if getattr(shp, "has_table", False):
                continue
            if getattr(shp, "has_text_frame", False):
                t = _text(shp)
                if keep_keywords and any(k in t for k in keep_keywords):
                    continue
                _remove_shape_xml(shp)
        except Exception:
            pass


def _clear_issue_content(slide) -> None:
    for shp in list(slide.shapes):
        try:
            if getattr(shp, "has_table", False) or getattr(shp, "has_text_frame", False):
                _remove_shape_xml(shp)
        except Exception:
            pass


def _add_textbox(slide, x, y, w, h, text, font_size, bold, color, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = _clean(text) or "—"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _set_cell_fill(cell, color: RGBColor) -> None:
    try:
        cell.fill.solid()
        cell.fill.fore_color.rgb = color
    except Exception:
        pass


def _set_cell(cell, text: str, font_size: int = 12, bold: bool | None = None, align=PP_ALIGN.LEFT, blank_ok: bool = False) -> None:
    cell.text = _clean(text) if blank_ok else (_clean(text) or "—")
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Pt(4)
    cell.margin_right = Pt(4)
    cell.margin_top = Pt(2)
    cell.margin_bottom = Pt(2)
    for p in cell.text_frame.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.name = "Microsoft YaHei"
            r.font.size = Pt(font_size)
            if bold is not None:
                r.font.bold = bold


def _split_text(text: str, limit: int) -> list[str]:
    text = _clean(text)
    if not text:
        return [""]
    out, buf = [], ""
    for part in re.split(r"\n\s*\n", text):
        part = part.strip()
        if not part:
            continue
        if len(buf) + len(part) + 2 <= limit:
            buf = (buf + "\n\n" + part).strip()
        else:
            if buf:
                out.append(buf)
            while len(part) > limit:
                out.append(part[:limit])
                part = part[limit:]
            buf = part
    if buf:
        out.append(buf)
    return out or [""]


def _meaningful_text(value: Any) -> bool:
    t = _clean(value)
    if not t or t in {"—", "-", "无", "NA", "N/A"}:
        return False
    if any(p in t for p in PLACEHOLDER_TEXTS):
        return False
    return len(re.sub(r"[\s\-—_：:|]+", "", t)) > 0


def _paginate_issue(issue: dict) -> list[dict]:
    basis = _clean(issue.get("basis", ""))
    desc = _clean(issue.get("description", ""))
    if not _meaningful_text(basis) and not _meaningful_text(desc):
        return []
    basis_parts = _split_text(basis, 430)
    desc_parts = _split_text(desc, 1150) if _meaningful_text(desc) else [""]
    total = max(len(basis_parts), len(desc_parts))
    pages = []
    for i in range(total):
        x = dict(issue)
        x["basis"] = basis_parts[i] if i < len(basis_parts) else ""
        x["description"] = desc_parts[i] if i < len(desc_parts) else ""
        x["title"] = ""
        x["_sub_page"] = i + 1
        x["_sub_total"] = total
        pages.append(x)
    return pages


def _has_issue_content(issue: dict) -> bool:
    return _meaningful_text(issue.get("basis", "")) or _meaningful_text(issue.get("description", ""))


def _estimate_lines(text: str, chars_per_line: int = 35) -> int:
    text = _clean(text)
    if not text:
        return 1
    parts = text.split("\n")
    return max(1, sum(max(1, math.ceil(len(p) / chars_per_line)) for p in parts))


def _render_cover(slide, context):
    _remove_text_shapes(slide)
    meta = context.get("meta", {})
    project = _clean(meta.get("project_name", "—"))
    center = _clean(meta.get("center_name", "—"))
    center_no = _clean(meta.get("center_no", ""))
    audit_date = _clean(meta.get("audit_date", "—"))
    audit_company = _clean(meta.get("audit_company", "北京万宁睿和医药科技有限公司"))

    title = f"{project}-{center}" + (f"（中心编号{center_no}）" if center_no else "")
    line_count = _estimate_lines(title, 35)
    title_h = min(1.45, max(0.80, 0.38 * line_count + 0.18))
    title_y = 3.22
    meeting_y = min(4.70, title_y + title_h + 0.12)
    time_y = meeting_y + 0.62
    company_y = time_y + 0.32

    _add_textbox(slide, 0.25, title_y, 12.70, title_h, title, 28, True, WHITE, PP_ALIGN.LEFT)
    _add_textbox(slide, 0.25, meeting_y, 12.70, 0.46, "中心稽查末次会议", 28, True, YELLOW, PP_ALIGN.LEFT)
    _add_textbox(slide, 0.35, time_y, 12.30, 0.30, f"时间：{audit_date}", 14, True, WHITE, PP_ALIGN.LEFT)
    _add_textbox(slide, 0.35, company_y, 12.30, 0.30, audit_company, 14, True, WHITE, PP_ALIGN.LEFT)


def _render_overview(slide, context):
    meta = context.get("meta", {})
    tables = _tables(slide)
    if not tables:
        return
    tbl = tables[0]
    rows = [
        [(0, 0, "方案名称", 16, True), (0, 1, _clean(meta.get("project_name", "—")), 16, False)],
        [(1, 0, "申办者", 16, True), (1, 1, _clean(meta.get("sponsor", "")), 16, False), (1, 2, "PI", 16, True), (1, 3, _clean(meta.get("pi", "—")), 16, False)],
        [(2, 0, "中心名称", 16, True), (2, 1, _clean(meta.get("center_name", "—")), 16, False), (2, 2, "中心入组\n情况", 16, True), (2, 3, _clean(meta.get("enrollment", "—")), 16, False)],
        [(3, 0, "稽查时间", 16, True), (3, 1, _clean(meta.get("audit_date", "—")), 16, False), (3, 2, "稽查员", 16, True), (3, 3, _clean(meta.get("auditor", "—")), 16, False)],
        [(4, 0, f"本次稽查{len(context.get('audited_subjects') or []) or 'x'}\n例受试者", 16, True), (4, 1, "、".join(context.get("audited_subjects") or ["—"]), 16, False)],
    ]
    for group in rows:
        for r, c, text, size, bold in group:
            if r < len(tbl.rows) and c < len(tbl.columns):
                _set_cell(tbl.cell(r, c), text, size, bold, PP_ALIGN.CENTER if c in (0, 2) else PP_ALIGN.LEFT, blank_ok=(r == 1 and c == 1))


def _render_counts(slide, context):
    counts = context.get("summary", {})
    tables = _tables(slide)
    for tbl, cats in zip(tables[:2], [LEFT_CATS, RIGHT_CATS]):
        _set_cell(tbl.cell(0, 0), "分类", 12, True, PP_ALIGN.CENTER)
        _set_cell(tbl.cell(0, 1), "数量", 12, True, PP_ALIGN.CENTER)
        for i, cat in enumerate(cats, start=1):
            if i >= len(tbl.rows):
                break
            _set_cell(tbl.cell(i, 0), cat, 12, False, PP_ALIGN.LEFT)
            _set_cell(tbl.cell(i, 1), str(counts.get(cat, 0) or "—"), 12, True, PP_ALIGN.CENTER)


def _render_issue(slide, category, issue, idx, total):
    _clear_issue_content(slide)
    sub_page = int(issue.get("_sub_page", 1))
    sub_total = int(issue.get("_sub_total", 1))
    tag = f"（{idx}/{total}）" if total > 1 else ""
    cont = f"  续{sub_page}/{sub_total}" if sub_total > 1 else ""
    page_title = f"问题分类：{category}{tag}{cont}"
    basis = _clean(issue.get("basis", ""))
    desc = _clean(issue.get("description", ""))
    _add_textbox(slide, 0.70, 0.52, 12.0, 0.45, page_title, 20, True, BLACK, PP_ALIGN.LEFT)
    rows = [("问题依据", basis), ("问题描述", desc)]
    shape = slide.shapes.add_table(2, 2, Inches(0.65), Inches(1.22), Inches(12.05), Inches(5.45))
    tbl = shape.table
    tbl.columns[0].width = Inches(1.05)
    tbl.columns[1].width = Inches(11.00)
    tbl.rows[0].height = Inches(1.55)
    tbl.rows[1].height = Inches(3.90)
    for r, (label, value) in enumerate(rows):
        _set_cell(tbl.cell(r, 0), label, 14, True, PP_ALIGN.CENTER)
        size = 11 if label == "问题描述" else 10
        _set_cell(tbl.cell(r, 1), value, size, False, PP_ALIGN.LEFT, blank_ok=(label == "问题依据"))


def _risk_score(issue: dict) -> int:
    category = issue.get("category", "其他")
    text = f"{issue.get('summary', '')}\n{issue.get('description', '')}\n{issue.get('basis', '')}\n{issue.get('full_text', '')}"
    score = RISK_CATEGORY_WEIGHT.get(category, 10)
    severity = _clean(issue.get("severity", ""))
    if severity in {"高", "重大", "严重", "Major", "High", "high", "major"}:
        score += 35
    elif severity in {"中", "一般", "Medium", "medium"}:
        score += 12
    for kw, weight in RISK_KEYWORDS.items():
        if kw.lower() in text.lower():
            score += weight
    score += min(12, len(text) // 160)
    return score


def _risk_advice(category: str, text: str) -> str:
    lower_text = text.lower()
    if "知情同意" in category or "icf" in lower_text:
        return "逐例复核ICF签署版本、签署日期/时间、签署人资质、授权委托及受试者权益告知证据；准备签署过程说明和更正/补充记录。"
    if "安全性" in category or "sae" in lower_text or "susar" in lower_text:
        return "逐例核对AE/SAE从原始病历到EDC及上报系统的完整链条，确认严重性、相关性、转归、上报时限和随访闭环证据。"
    if "方案依从" in category or "入排" in text or "方案偏离" in text:
        return "围绕入排标准、访视窗口、给药/检查流程和偏离记录建立逐例核查清单，提前形成偏离判定、医学解释和CAPA闭环材料。"
    if "原始文件" in category or "HIS" in category or "LIS" in category or "PACS" in category:
        return "提前完成原始病历、HIS/LIS/PACS、源文件与EDC一致性复核，标记差异原因，准备可追溯证据和研究者确认说明。"
    if "CRF" in category or "EDC" in lower_text:
        return "导出EDC关键字段核查清单，重点复核录入及时性、逻辑一致性、Query关闭证据及源数据支持。"
    if "试验用药" in category:
        return "复核药品接收、储存温度、发放、回收、清点、销毁及授权人员记录，确保账物卡一致并可追溯。"
    if "生物样本" in category:
        return "复核样本采集、处理、保存、运输、交接、检测结果回传全链条，重点确认时间窗、标签、温控和偏差处理记录。"
    if "伦理" in category:
        return "核对伦理批件、递交材料、方案/ICF版本、持续审查和安全性信息递交记录，确保版本和执行时间一致。"
    return "按问题清单逐项准备原始证据、研究者说明、整改记录和CAPA闭环材料；对同类问题开展横向复核，避免核查现场重复暴露。"


def _extract_top5_risks(context: dict) -> list[dict]:
    issues = [x for x in context.get("issues", []) if _has_issue_content(x)]
    enriched = []
    seen = set()
    for issue in issues:
        category = issue.get("category", "其他")
        desc = _clean(issue.get("description", "")) or _clean(issue.get("summary", ""))
        basis = _clean(issue.get("basis", ""))
        key = (category, _shorten(desc, 60))
        if key in seen:
            continue
        seen.add(key)
        full_text = f"{desc}\n{basis}"
        enriched.append({
            "category": category,
            "risk": _shorten(desc, 96),
            "reason": _shorten(basis if _meaningful_text(basis) else f"该问题涉及{category}，可能影响核查对试验质量和合规性的判断。", 86),
            "advice": _risk_advice(category, full_text),
            "score": _risk_score(issue),
        })
    enriched.sort(key=lambda x: x["score"], reverse=True)
    return enriched[:5]


def _render_risk_summary(slide, context):
    _clear_issue_content(slide)
    risks = _extract_top5_risks(context)
    _add_textbox(slide, 0.55, 0.38, 12.20, 0.48, "核查准备重点关注问题及迎检建议", 24, True, BLACK, PP_ALIGN.LEFT)
    if not risks:
        _add_textbox(slide, 0.75, 1.35, 11.80, 0.80, "本次上传文件中未识别到可用于提炼TOP5的问题内容，请复核Excel问题分类、问题描述和依据列是否完整。", 16, False, BLACK, PP_ALIGN.LEFT)
        return

    shape = slide.shapes.add_table(6, 4, Inches(0.45), Inches(1.10), Inches(12.45), Inches(5.70))
    tbl = shape.table
    widths = [0.55, 2.25, 4.45, 5.20]
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)
    tbl.rows[0].height = Inches(0.45)
    for r in range(1, 6):
        tbl.rows[r].height = Inches(1.05)
    headers = ["序号", "风险类别", "TOP高风险问题", "迎检建议"]
    for c, h in enumerate(headers):
        _set_cell(tbl.cell(0, c), h, 12, True, PP_ALIGN.CENTER)
        _set_cell_fill(tbl.cell(0, c), LIGHT_BLUE)
    for r in range(1, 6):
        if r <= len(risks):
            item = risks[r - 1]
            values = [str(r), item["category"], f"{item['risk']}\n依据/风险逻辑：{item['reason']}", item["advice"]]
        else:
            values = [str(r), "—", "—", "—"]
        for c, v in enumerate(values):
            size = 10 if c >= 2 else 11
            align = PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT
            _set_cell(tbl.cell(r, c), v, size, c == 0, align)
            if c == 0:
                _set_cell_fill(tbl.cell(r, c), LIGHT_YELLOW)
    _add_textbox(slide, 0.55, 6.92, 12.25, 0.25, "注：本页基于本次稽查发现自动提炼，用于核查准备优先级排序；正式迎检材料需结合项目医学判断及原始证据人工复核。", 9, False, GRAY, PP_ALIGN.LEFT)


def _copy_if_exists(prs, src_no):
    if src_no and 1 <= src_no <= len(prs.slides):
        return _copy_slide(prs, src_no)
    return None


def _find_slide_no_by_keyword(prs: Presentation, keyword: str, default_no: int | None = None) -> int | None:
    for idx, slide in enumerate(prs.slides, start=1):
        if keyword in _slide_text(slide):
            return idx
    return default_no


def render_ppt(context, output_path, template_path=None):
    template = Path(template_path or DEFAULT_TEMPLATE_PATH)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    if not template.exists() or template.stat().st_size < 1024 * 100:
        raise FileNotFoundError(f"未找到有效PPT模板：{template}。请在页面先上传你的稽查总结会模板。")

    prs = Presentation(str(template))
    original_count = len(prs.slides)
    if original_count < SLIDE_COUNTS:
        raise RuntimeError(f"模板页数不足：当前 {original_count} 页，至少需要 {SLIDE_COUNTS} 页。请上传完整的稽查总结会模板。")

    summary_slide_no = _find_slide_no_by_keyword(prs, "核查准备重点关注问题及迎检建议", original_count - 2 if original_count >= 3 else None)
    qa_slide_no = original_count - 1 if original_count >= 2 else None
    ending_slide_no = original_count if original_count >= 1 else None
    issue_template_no = SLIDE_ISSUE if original_count >= SLIDE_ISSUE else SLIDE_COUNTS

    slide = _copy_slide(prs, SLIDE_COVER)
    _render_cover(slide, context)
    _copy_slide(prs, SLIDE_THANKS)
    _copy_slide(prs, SLIDE_TOC)
    _copy_slide(prs, SLIDE_PART1)
    slide = _copy_slide(prs, SLIDE_OVERVIEW)
    _render_overview(slide, context)
    _copy_slide(prs, SLIDE_SCOPE)
    _copy_if_exists(prs, SLIDE_COMMON)
    _copy_if_exists(prs, SLIDE_INDIVIDUAL)
    _copy_slide(prs, SLIDE_PART2)
    slide = _copy_slide(prs, SLIDE_COUNTS)
    _render_counts(slide, context)

    for cat in context.get("standard_categories", []):
        cat_issues = [x for x in context.get("issues", []) if x.get("category") == cat and _has_issue_content(x)]
        if not cat_issues:
            continue
        for i, issue in enumerate(cat_issues, start=1):
            for page_issue in _paginate_issue(issue):
                if not _has_issue_content(page_issue):
                    continue
                slide = _copy_slide(prs, issue_template_no)
                _render_issue(slide, cat, page_issue, i, len(cat_issues))

    if summary_slide_no:
        slide = _copy_slide(prs, summary_slide_no)
        _render_risk_summary(slide, context)
    _copy_if_exists(prs, qa_slide_no)
    _copy_if_exists(prs, ending_slide_no)

    _remove_original_template_slides(prs, original_count)
    _remove_template_empty_slides(prs)
    prs.save(str(output_path))
    return output_path
